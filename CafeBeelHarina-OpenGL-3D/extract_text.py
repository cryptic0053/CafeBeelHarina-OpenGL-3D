import os
import glob
from pptx import Presentation
import PyPDF2
from collections import Counter
import re

def extract_pptx(filepath):
    text = []
    try:
        prs = Presentation(filepath)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    filtered_text = shape.text.strip()
                    if len(filtered_text) > 10: # Only meaningful text
                        text.append(filtered_text)
    except Exception as e:
        print(f"Error reading {filepath}: {e}")
    return "\n".join(text)

def extract_pdf(filepath):
    text = []
    try:
        with open(filepath, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text.append(page_text.strip())
    except Exception as e:
        print(f"Error reading {filepath}: {e}")
    return "\n".join(text)

all_text = ""
print("Extracting PPTX files...")
for f in glob.glob("*.pptx"):
    all_text += extract_pptx(f) + "\n"

print("Extracting PDF files...")
for f in glob.glob("*.pdf"):
    all_text += extract_pdf(f) + "\n"

with open("extracted_text.txt", "w", encoding="utf-8") as out:
    out.write(all_text)

# Let's do a quick naive summary by finding the most common meaningful words
words = re.findall(r'\b[A-Za-z-]{5,}\b', all_text.lower())
stopwords = {"which", "their", "there", "about", "would", "these", "other", "where", "after", "could", "should", "using", "between", "through", "value", "given", "example"}
filtered_words = [w for w in words if w not in stopwords]
counts = Counter(filtered_words).most_common(100)

print("Top 100 keywords found across all slides:")
for word, count in counts:
    print(f"{word}: {count}")

print(f"Extraction complete. Total characters extracted: {len(all_text)}")
